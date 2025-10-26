import os
import tempfile
import requests

import pandas as pd
from PIL import Image
import pytesseract
from pptx import Presentation
import shutil
from langdetect import detect

from fastapi import HTTPException

from langchain_community.document_loaders import PyMuPDFLoader, Docx2txtLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter
from  langchain_text_splitters.sentence_transformers import SentenceTransformersTokenTextSplitter  # give better results but slow can use later for project
from langchain.schema import Document
from langchain_google_genai import ChatGoogleGenerativeAI


# directory for saved model for SentenceTransformersTokenTextSplitter
MODEL_DIR = os.path.join("/tmp", "e5-large-v2")

# for storing chunk (saving  timeeeeeeeeeeee)
chunk_dict= {}

GOOGLE_API_KEY2 = os.getenv("gemini_api_key2")
llm = ChatGoogleGenerativeAI(
                               model="gemini-2.0-flash",
                               api_key=GOOGLE_API_KEY2,
                           )
    
def load_excel(path: str) -> list[Document]:
    dfs = pd.read_excel(path, sheet_name=None)
    docs = []
    for sheet_name, df in dfs.items():
        text = df.to_csv(index=False)
        docs.append(Document(page_content=text, metadata={"sheet": sheet_name}))
    return docs


""" I am not using zip further as no questions were asked and also the nested zip has no data .
     i did not able to submit this coded approaches for differnet types of files due to sudden change of level.
     do not delete it  , may be all levels will reopened on last day"""


def load_zip(path: str, depth: int = 0, base_dir="/tmp/unzipped") -> list[Document]:
    extracted_docs = []
    extract_dir = os.path.join(base_dir, f"level_{depth}")
    os.makedirs(extract_dir, exist_ok=True)

    with zipfile.ZipFile(path, 'r') as archive:
        archive.extractall(extract_dir)

    for name in os.listdir(extract_dir):
        file_path = os.path.join(extract_dir, name)
        
        if name.endswith(".zip"):
            extracted_docs.extend(load_zip(file_path, depth + 1, base_dir))  # Recursive call
        elif name.endswith(".pdf"):
            loader = PyMuPDFLoader(file_path)
            extracted_docs += loader.load()
        elif name.endswith(".docx"):
            loader = Docx2txtLoader(file_path)
            extracted_docs += loader.load()
        elif name.endswith(".txt"):
            with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
                extracted_docs.append(Document(page_content=f.read()))
        elif name.endswith((".png", ".jpg", ".jpeg")):
            image = Image.open(file_path)
            text = pytesseract.image_to_string(image)
            extracted_docs.append(Document(page_content=text))

    return extracted_docs


def load_image(path: str) -> list[Document]:
    image = Image.open(path)
    text = pytesseract.image_to_string(image)
    return [Document(page_content=text)]

def load_pptx(path: str) -> list[Document]:
    prs = Presentation(path)
    full_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                full_text.append(shape.text)
            elif shape.shape_type == 13 and shape.image:  # PICTURE shape
                image = shape.image.blob
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as img_tmp:
                    img_tmp.write(image)
                    img_path = img_tmp.name
                try:
                    img_text = pytesseract.image_to_string(Image.open(img_path))
                    if img_text.strip():
                        full_text.append(img_text.strip())
                finally:
                    os.remove(img_path)
    return [Document(page_content="\n".join(full_text))]

def load_and_chunk(url: str) -> list[Document]:
    print(url)
       
    if url not in chunk_dict:
        print("processing new url")
        resp = requests.get(url)
        if resp.status_code != 200:
            raise HTTPException(400, "Could not download document")

        content_type = resp.headers.get("Content-Type", "").lower()
        url_lower = url.lower()

        try:
            if "application/pdf" in content_type or ".pdf" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                loader = PyMuPDFLoader(tmp_path)
                docs = loader.load_and_split()

            elif "application/vnd.openxmlformats-officedocument.wordprocessingml.document" in content_type or ".docx" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".docx") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                loader = Docx2txtLoader(tmp_path)
                docs = loader.load_and_split()

            elif "text/plain" in content_type or ".txt" in url_lower:
                text = resp.content.decode("utf-8", errors="ignore")
                docs = [Document(page_content=text)]

            elif ".xlsx" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".xlsx") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                docs = load_excel(tmp_path)

            elif ".zip" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".zip") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                text = "empty file"    
                docs = [Document(page_content=text)]

            elif ".png" in url_lower or ".jpg" in url_lower or ".jpeg" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".png") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                docs = load_image(tmp_path)

            elif ".pptx" in url_lower:
                with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                    tmp.write(resp.content)
                    tmp_path = tmp.name
                docs = load_pptx(tmp_path)

            else:
                raise HTTPException(400, f"Unsupported document type: {content_type}")

        finally:
            if 'tmp_path' in locals() and os.path.exists(tmp_path):
                os.remove(tmp_path)

        
        
        min_len = 100 if len(docs[0].page_content) >= 100 else len(docs[0].page_content)
        lang = detect(docs[0].page_content[:min_len])

        text = ""
        if lang != "en":
            for doc in docs:
                translated = llm.invoke(f"Translate this into English:\n{doc.page_content}")
                text = text + translated.content

        else:
            full_text = "\n".join([doc.page_content for doc in docs])
            
        splitter = SentenceTransformersTokenTextSplitter(
            model_name=MODEL_DIR, tokens_per_chunk=512, chunk_overlap=90
        )
        
        if lang !="en":
            chunk_dict[url] = splitter.create_documents([text])
        else:    
            chunk_dict[url] = splitter.create_documents([full_text])
            
        return chunk_dict[url]
    else:
        print("stored chunk")
        return chunk_dict[url]
